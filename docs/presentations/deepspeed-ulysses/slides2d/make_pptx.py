import glob, os
from pptx import Presentation
from pptx.util import Inches
HERE = os.path.dirname(os.path.abspath(__file__))
FR = os.path.join(HERE, 'frames')
OUT = os.path.join(HERE, '..', 'deepspeed-ulysses.pptx')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank
pngs = sorted(glob.glob(os.path.join(FR, '*.png')))
for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(OUT)
print('wrote', os.path.normpath(OUT), '·', len(pngs), 'slides')
